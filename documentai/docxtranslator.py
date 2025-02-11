import json
import pathlib
import sys

from pptx import Presentation
from docx import Document
from pptx.util import Pt
from langchain_community.document_loaders import UnstructuredWordDocumentLoader


class DOCXPPTHindi2English:

    def __init__(self, file_name, fn, docx=False):
        self.fn = fn
        self.file_name = file_name
        if docx:
            self.prs = None
            self.docx = Document(file_name)
        else:
            self.prs = Presentation(file_name)
        if pathlib.Path(file_name + '_tmp').exists():
            print(file_name+'_tmp')
            with open( file_name + '_tmp', 'r') as f:
                try:
                   self.storage_dict = json.load(f)
                   f.close()
                except:
                    raise Exception("prior session does not exist. Set prev_session to False and run again")
        else:
            self.storage_dict = {}
    def serialize_storage_dict(self):
        with open(self.file_name + '_tmp', 'w') as f:
            json.dump(self.storage_dict,f)

    def import_dict_to_word(self, save_as):
       print(self.storage_dict.keys())
       i = 0
       for paragraph in self.docx.paragraphs:
           if str(i) in self.storage_dict:
               paragraph.text = self.storage_dict[str(i)]
               print('replacing', end='.')
           i += 1
           print( i, end='.')

       for tbl in self.docx.tables:
           for col in tbl.columns:
               for cell in col.cells:
                   for paragraph in cell.paragraphs:
                       if str(i) in self.storage_dict:
                           paragraph.text = self.storage_dict[str(i)]
                           print('replacing', end='.')

                       i = i + 1

       self.docx.save(save_as)

    def convert_word(self, save_as=None, save_freq = 10,starting=0):
        i = 0
        for paragraph in self.docx.paragraphs:
            if i < starting: continue
            paragraph.text = self.fn(paragraph.text)
            i += 1
            print('\r',i,'.')
            if (i % save_freq) == 0: self.docx.save(save_as)

        for tbl in self.docx.tables:
            for col in tbl.columns:
                for cell in col.cells:
                    for paragraph in cell.paragraphs:
                            if i < starting: continue
                            paragraph.text = self.fn(paragraph.text)
                            i += 1
                            print('\r', i, '.')
                            if save_as:
                                if (i % save_freq) == 0:self.docx.save(save_as)

        if save_as:
             self.docx.save(save_as)
    def convert_word_intermediate(self, save_as=None, save_freq = 10,starting=0):
        i = 0
        total = len(self.docx.paragraphs)
        for tbl in self.docx.tables:
            for col in tbl.columns:
                for cell in col.cells:
                    total += len(cell.paragraphs)
        print("total paragraphs to be translated:", total)
        for paragraph in self.docx.paragraphs:
            if i < starting: continue
            if str(i) in self.storage_dict: continue
            self.storage_dict[str(i)] = self.fn(paragraph.text)
            i += 1
            print(i, f'/{total}', end='.')
            if (i % save_freq) == 0: self.serialize_storage_dict()

        for tbl in self.docx.tables:
            for col in tbl.columns:
                for cell in col.cells:
                    for paragraph in cell.paragraphs:
                            if i < starting: continue
                            if str(i) in self.storage_dict: continue
                            self.storage_dict[str(i)] = self.fn(paragraph.text)
                            i += 1
                            print( i,f'/{total}', end='.')
                            if save_as:
                                if (i % save_freq) == 0:self.serialize_storage_dict()

        if save_as:
             self.serialize_storage_dict()


    def convert(self, save_as=None):
        i = 0
        for slide in self.prs.slides:
            i += 1
            #print(f'slide no {i}')
            for shape in slide.shapes:
                if shape.has_table:
                    #print("table found")
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                paragraph.text = self.fn(paragraph.text)
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    paragraph.text = self.fn(paragraph.text)
        if save_as: self.prs.save(save_as)

    def save(self, file_name):
        if file_name: self.prs.save(file_name)





class LangChainDocx:
    def __init__(self, path):
        loader = UnstructuredWordDocumentLoader(path)
        self.data = loader.load()
if __name__ == '__main__':
    from hindi2englishml import hindi2english
    #word_path = '/Users/ranvirprasad/Downloads/02.docx'
    #lc = LangChainDocx(word_path)
    #print(lc.data)

    file_name = '/Users/ranvirprasad/Downloads/Ch_09_Industry.docx'
    converter = DOCXPPTHindi2English(file_name, docx=True, fn=hindi2english)
    #converter.convert_word_intermediate(save_as=True, starting=0)
    converter.import_dict_to_word('/Users/ranvirprasad/Downloads/temp.docx')
    #converter.convert_word(save_as='All_Chap_english.docx', starting=61)

    #converter1 = PPTKruti2Unicode('/Users/ranvirprasad/Downloads/Rahul Gahlot Shikayat Aakhya New - Final.docx', docx=True)
    #converter1.convert_word('converted_word2.docx')
